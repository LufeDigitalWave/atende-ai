"""Generate WhatsApp Agents pitch deck PPTX."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colors
NAVY = RGBColor(0x1E, 0x27, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
VIOLET = RGBColor(0xA8, 0x55, 0xF7)
GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK = RGBColor(0x0A, 0x0F, 0x1F)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# --- Slide 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
         "Atende AI", font_size=48, bold=True, color=WHITE)
add_text(slide, Inches(1), Inches(3.3), Inches(11), Inches(1),
         "Agentes de IA para WhatsApp que respondem, qualificam e encaminham leads 24/7.",
         font_size=22, color=CYAN)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Lufe Digital Wave", font_size=14, color=GRAY)

# --- Slide 2: Problema ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "O problema", font_size=32, bold=True, color=VIOLET)
add_bullet_list(slide, Inches(1), Inches(2), Inches(10), Inches(4), [
    "Lead manda mensagem no WhatsApp e ninguem responde a tempo.",
    "Fora do horario comercial, o lead esfria ou vai para o concorrente.",
    "Vendedor recebe contato sem contexto: nao sabe o que o cliente quer.",
    "FAQ basico consome tempo da equipe humana todos os dias.",
    "Informacoes espalhadas entre planilhas, CRM e cabeca do atendente.",
], font_size=16)

# --- Slide 3: Solucao ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "A solucao: WhatsApp Agents", font_size=32, bold=True, color=CYAN)
add_bullet_list(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), [
    "Agente de IA 24/7 no WhatsApp da empresa.",
    "Responde instantaneamente com linguagem natural.",
    "Qualifica leads e atualiza CRM em tempo real.",
    "Consulta base de conhecimento (precos, FAQ, politicas).",
    "Encaminha para humano com contexto completo quando precisa.",
    "Usa API oficial Meta: sem risco de banimento.",
], font_size=16)

# --- Slide 4: Demo ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Demo publica", font_size=32, bold=True, color=WHITE)
add_text(slide, Inches(1), Inches(2), Inches(10), Inches(1),
         "https://atendeai.lufedigitalwave.com.br", font_size=20, color=CYAN)
add_bullet_list(slide, Inches(1), Inches(3.2), Inches(10), Inches(3.5), [
    "Escolha um nicho e veja a IA conversando em tempo real.",
    "CRM ao vivo preenche enquanto o lead conversa.",
    "Score e funil atualizam automaticamente.",
    "30+ nichos suportados pela Factory v3.",
    "Dados 100% ficticios no portfolio. A IA e real.",
], font_size=16)

# --- Slide 5: Catalogo de Agentes ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "6 agentes, uma base tecnica", font_size=32, bold=True, color=VIOLET)
agents_text = [
    "1. SDR Agent — vendas e qualificacao de leads.",
    "2. Support Agent — atendimento e suporte.",
    "3. Appointment Agent — agendamento automatico.",
    "4. FAQ/RAG Agent — base de conhecimento.",
    "5. Civic Agent — atendimento publico/protocolos.",
    "6. Collections Agent — follow-up e reativacao.",
]
add_bullet_list(slide, Inches(1), Inches(2), Inches(10), Inches(4.5), agents_text, font_size=16)

# --- Slide 6: Stack ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Stack e seguranca", font_size=32, bold=True, color=CYAN)
add_bullet_list(slide, Inches(1), Inches(2), Inches(5), Inches(4.5), [
    "FastAPI + React + PostgreSQL + pgvector",
    "WhatsApp Cloud API (Meta oficial)",
    "LLM multi-provider (Claude, OpenAI)",
    "Docker Compose + VPS",
    "CI GitHub Actions",
], font_size=16)
add_bullet_list(slide, Inches(6.5), Inches(2), Inches(5), Inches(4.5), [
    "Rate limit + budget diario",
    "Kill switch operacional",
    "Session TTL + soft delete",
    "PII sanitizer em logs",
    "Allowlist WhatsApp em contatos",
], font_size=16)

# --- Slide 7: Pacotes ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
         "Pacotes", font_size=32, bold=True, color=WHITE)

packages = [
    ("Starter", "R$ 1.500-3.000", "R$ 300-800/mes", "1 agente, handoff, 7-10 dias"),
    ("Pro", "R$ 3.500-8.000", "R$ 800-2.000/mes", "RAG, CRM/agenda, metricas"),
    ("Business", "R$ 8.000+", "R$ 2.000+/mes", "Multiagente, integracoes, operacao"),
]

for i, (name, impl, mens, desc) in enumerate(packages):
    y = Inches(2.2 + i * 1.6)
    add_text(slide, Inches(1), y, Inches(3), Inches(0.5), name, font_size=20, bold=True, color=CYAN)
    add_text(slide, Inches(1), y + Inches(0.5), Inches(4), Inches(0.4), f"Implantacao: {impl}", font_size=14, color=WHITE)
    add_text(slide, Inches(5), y + Inches(0.5), Inches(4), Inches(0.4), f"Mensalidade: {mens}", font_size=14, color=WHITE)
    add_text(slide, Inches(1), y + Inches(0.9), Inches(10), Inches(0.4), desc, font_size=13, color=GRAY)

# --- Slide 8: CTA ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
         "Proximo passo", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
         "Posso montar um piloto rapido com seu nicho para voce ver funcionando antes de fechar.",
         font_size=20, color=CYAN, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "Luiz Felipe — Tech Lead IA / Especialista WhatsApp Cloud API",
         font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "https://atendeai.lufedigitalwave.com.br",
         font_size=14, color=CYAN, alignment=PP_ALIGN.CENTER)

# Save
output = "C:/tmp/atende-ai/docs/commercial/whatsapp-agents-pitch.pptx"
prs.save(output)
print(f"Saved: {output}")
